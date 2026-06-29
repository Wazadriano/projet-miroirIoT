# -*- coding: utf-8 -*-
"""
Generateur du PowerPoint de soutenance Smart Mirror KBEAUTY (RNCP 37046).

Source unique et DYNAMIQUE : docs/PRESENTATION-STORYBOARD.md.
Le script parse le markdown du storyboard et genere une slide par bloc
"### Slide N - ...". Le contenu n'est plus code en dur : titres, puces et
notes orateur sont extraits du markdown, ce qui garantit que le .pptx suit
toujours la derniere version du storyboard.

Regle : une slide par slide du storyboard, contenu visible MINIMAL (puces du
storyboard), notes orateur completes (Notes orateur + Note interne concatenees),
identite visuelle prune/or/ivoire (sandwich aubergine pour ouverture, accroche,
transition, section et cloture).

Aucune image n'est telechargee : chaque slide reserve un espace visuel et
indique en petit la suggestion d'icone/illustration prevue ("Visuel suggere").
"""

import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Palette (identite visuelle du storyboard, section 1)
# ---------------------------------------------------------------------------
PLUM = RGBColor(0x6B, 0x2D, 0x5C)        # primaire
GOLD = RGBColor(0x9C, 0x7A, 0x3C)        # secondaire / accents / chiffres
IVORY = RGBColor(0xF7, 0xF3, 0xEE)       # fond clair
AUBERGINE = RGBColor(0x2A, 0x13, 0x26)   # fond sombre (sandwich)
ANTHRACITE = RGBColor(0x22, 0x20, 0x2A)  # texte sur fond clair
TERRACOTTA = RGBColor(0xB5, 0x56, 0x2F)  # alerte / honnetete / RGPD
SAGE = RGBColor(0x5E, 0x7A, 0x5A)        # validation / MVP / tests verts
LIGHT_ON_DARK = RGBColor(0xF2, 0xE9, 0xDF)
SOFT_VISUAL = RGBColor(0xEC, 0xE3, 0xD7)  # cadre placeholder visuel
GREY = RGBColor(0x8A, 0x84, 0x80)

TITLE_FONT = "Georgia"        # serif (esprit institut premium)
BODY_FONT = "Calibri"         # sans-serif lisible

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FOOTER_TEXT = "Adriano - RNCP 37046 - Smart Mirror KBEAUTY - 2026"

HERE = os.path.dirname(os.path.abspath(__file__))
STORYBOARD = os.path.normpath(os.path.join(HERE, "..", "PRESENTATION-STORYBOARD.md"))
OUTPUT = os.path.normpath(os.path.join(HERE, "..", "Soutenance-KBEAUTY-RNCP.pptx"))

# Mots-cles d'en-tete qui declenchent une slide "sandwich" (fond sombre, titre
# centre) : ouverture, accroche, transition, slide de section, remerciements.
SANDWICH_KEYWORDS = ("ouverture", "accroche", "transition", "remerciement", "section")

# Mapping numero de slide -> (section affichee, couleur du bandeau).
# Derive de la table de minutage et de la table des codes par section du
# storyboard (section 1 et 2). C'est de la metadonnee visuelle, pas du contenu.
SECTION_RANGES = [
    (1, 2, "Ouverture", AUBERGINE),
    (3, 6, "Qui je suis", PLUM),
    (7, 8, "Entreprise et client", PLUM),
    (9, 13, "Probleme et marche", PLUM),
    (14, 16, "Cahier des charges", PLUM),
    (17, 17, "Devis", PLUM),
    (18, 19, "Gestion de projet", PLUM),
    (20, 28, "Conception et technique", AUBERGINE),
    (29, 29, "Demo et resultats", SAGE),
    (30, 31, "Bilan et perspectives", PLUM),
    (32, 32, "Remerciements", AUBERGINE),
]


def section_for(n):
    for lo, hi, name, banner in SECTION_RANGES:
        if lo <= n <= hi:
            return name, banner
    return "Soutenance", PLUM


# ---------------------------------------------------------------------------
# Parsing du storyboard markdown
# ---------------------------------------------------------------------------

FIELD_RE = re.compile(r"^\*\*\s*(.+?)\s*:\*\*\s*(.*)$")
SLIDE_HEADER_RE = re.compile(r"^###\s+Slide\s+(\d+)\s*-\s*(.*)$", re.MULTILINE)


def _clean_inline(text):
    """Nettoie une portion de texte : retire le gras markdown et espaces."""
    text = text.replace("**", "").strip()
    return text


def _classify(label):
    low = label.lower()
    if low == "titre":
        return "title"
    if low.startswith("contenu visible"):
        return "bullets"
    if low.startswith("notes orateur"):
        return "notes"
    if low.startswith("note interne"):
        return "notes"
    if low.startswith("visuel"):
        return "visual"
    return None


def _section_bounds(text, start_marker, end_markers):
    start = text.find(start_marker)
    if start == -1:
        return None
    start += len(start_marker)
    end = len(text)
    for marker in end_markers:
        idx = text.find(marker, start)
        if idx != -1:
            end = min(end, idx)
    return text[start:end]


def _parse_slide_block(num, header, body):
    """Parse le corps d'un bloc slide en champs structures."""
    fields = {"title": None, "bullets_raw": [], "visual": None}
    notes = []  # liste de (label, text)
    current_type = None
    current_lines = []
    current_label = None

    def flush():
        if current_type is None:
            return
        joined = "\n".join(current_lines).strip()
        if current_type == "title":
            if joined:
                fields["title"] = _clean_inline(joined)
        elif current_type == "visual":
            if joined:
                fields["visual"] = _clean_inline(joined)
        elif current_type == "bullets":
            fields["bullets_raw"].extend(current_lines)
        elif current_type == "notes":
            if joined:
                notes.append((current_label, joined))

    for line in body.splitlines():
        m = FIELD_RE.match(line.strip())
        if m:
            flush()
            label, inline = m.group(1), m.group(2)
            current_type = _classify(label)
            current_label = label
            current_lines = [inline] if inline.strip() else []
        else:
            if current_type is not None:
                current_lines.append(line)
    flush()

    # Puces : on garde les lignes commencant par "-", on aplatit l'imbrication
    # et on retire les lignes d'intro qui se terminent par ":".
    bullets = []
    for raw in fields["bullets_raw"]:
        stripped = raw.strip()
        if not stripped.startswith("-"):
            continue
        text = _clean_inline(stripped.lstrip("-").strip())
        if not text:
            continue
        if text.endswith(":"):
            continue
        bullets.append(text)

    # Notes : concatenation Notes orateur + Note interne (ordre du document).
    note_chunks = []
    for label, txt in notes:
        if label.strip().lower() == "notes orateur":
            note_chunks.append(txt)
        else:
            note_chunks.append("[" + _clean_inline(label) + "]\n" + txt)
    notes_text = "\n\n".join(note_chunks).strip()

    title = fields["title"] or _clean_inline(header)
    section, banner = section_for(num)
    dark = any(k in header.lower() for k in SANDWICH_KEYWORDS)
    if dark:
        banner = AUBERGINE
    todo = "a completer par adriano" in notes_text.lower()

    return dict(
        n=num,
        header=_clean_inline(header),
        section=section,
        banner=banner,
        dark=dark,
        title=title,
        bullets=bullets,
        visual=fields["visual"] or "",
        notes=notes_text,
        todo=todo,
    )


def parse_storyboard(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()

    main = _section_bounds(
        text, "## 3. Storyboard slide par slide",
        ["## 4.", "## 5."],
    )
    if main is None:
        raise RuntimeError("Section '## 3. Storyboard slide par slide' introuvable")

    # Decoupage robuste sur les en-tetes "### Slide N - ...".
    headers = list(SLIDE_HEADER_RE.finditer(main))

    slides = []
    for i, m in enumerate(headers):
        num = int(m.group(1))
        header = m.group(2).strip()
        body_start = m.end()
        body_end = headers[i + 1].start() if i + 1 < len(headers) else len(main)
        body = main[body_start:body_end]
        slides.append(_parse_slide_block(num, header, body))

    # --- Slides annexes (backup) ---
    backup_raw = _section_bounds(text, "## 4.", ["## 5."])
    backups = []
    if backup_raw:
        for line in backup_raw.splitlines():
            bm = re.match(r"^-\s*(A\d+)\s*:\s*(.+)$", line.strip())
            if bm:
                code = bm.group(1)
                desc = _clean_inline(bm.group(2)).rstrip(".")
                first = re.split(r"\s*[\(:]", desc)[0].strip()
                backups.append(dict(
                    code=code,
                    title="Annexe " + code + " - " + first,
                    bullets=[desc],
                    visual="Detail a presenter a l'oral si question du jury",
                    notes=desc,
                ))

    return slides, backups


# ---------------------------------------------------------------------------
# Helpers de rendu
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color, font=BODY_FONT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.color.rgb = color
    f.name = font
    f.bold = bold
    f.italic = italic
    return tb


def add_bullets(slide, x, y, w, h, items, size, color, font=BODY_FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.05
        # puce or
        rb = p.add_run()
        rb.text = "▪  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = GOLD
        rb.font.name = font
        rb.font.bold = True
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = font
    return tb


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text or ""


def add_footer(slide, on_dark):
    color = RGBColor(0xC9, 0xB8, 0x8C) if on_dark else GREY
    add_text(slide, Inches(0.5), Inches(7.06), Inches(9.5), Inches(0.35),
             FOOTER_TEXT, 10, color, font=BODY_FONT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)


def add_pagenum(slide, label, on_dark):
    color = GOLD
    add_text(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
             label, 12, color, font=BODY_FONT, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_visual_placeholder(slide, x, y, w, h, text, on_dark):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x3A, 0x21, 0x35) if on_dark else SOFT_VISUAL
    box.line.color.rgb = GOLD
    box.line.width = Pt(1)
    box.line.dash_style = 2  # dashed
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Visuel suggere : " + text if text else "Visuel suggere"
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.name = BODY_FONT
    r.font.color.rgb = RGBColor(0xC9, 0xB8, 0x8C) if on_dark else GOLD
    return box


def add_todo_badge(slide, x, y):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                                   Inches(3.7), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TERRACOTTA
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[A COMPLETER PAR ADRIANO]"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = BODY_FONT
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _eyebrow(section, header):
    """Bandeau de section + rappel du titre d'en-tete du storyboard."""
    base = section.upper()
    if header and header.strip().lower() != section.strip().lower():
        return base + "   -   " + header
    return base


# ---------------------------------------------------------------------------
# Construction de la presentation
# ---------------------------------------------------------------------------

def build():
    slides_data, backups = parse_storyboard(STORYBOARD)
    total = len(slides_data)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for s in slides_data:
        slide = prs.slides.add_slide(blank)
        dark = s["dark"]
        banner = s["banner"]
        on_dark = dark

        # Fond
        set_bg(slide, AUBERGINE if dark else IVORY)

        # Bandeau superieur de section + filet or
        add_rect(slide, 0, 0, SLIDE_W, Inches(0.22), banner)
        add_rect(slide, 0, Inches(0.22), SLIDE_W, Emu(28000), GOLD)

        # Eyebrow (section + en-tete du storyboard)
        sec_color = GOLD if not dark else RGBColor(0xC9, 0xB8, 0x8C)
        add_text(slide, Inches(0.5), Inches(0.32), Inches(11.0), Inches(0.4),
                 _eyebrow(s["section"], s["header"]), 12, sec_color,
                 font=BODY_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        title_color = GOLD if dark else PLUM

        if dark:
            # Slide sandwich : titre centre, phrase forte, puces en ligne
            center = None
            rest = list(s["bullets"])
            if rest:
                center = rest.pop(0)
                m = re.match(
                    r'^(Sous-titre|Une phrase centrale|Une phrase de bascule|Une ligne)\s*:\s*(.+)$',
                    center)
                if m:
                    center = m.group(2)
                center = center.strip().strip('“”"').rstrip('.').strip('“”"')

            add_text(slide, Inches(1.0), Inches(2.0), Inches(11.33), Inches(1.6),
                     s["title"], 48, title_color, font=TITLE_FONT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_rect(slide, Inches(5.16), Inches(3.7), Inches(3.0), Emu(28000), GOLD)
            if center:
                add_text(slide, Inches(1.2), Inches(3.9), Inches(10.93), Inches(1.0),
                         center, 24, LIGHT_ON_DARK, font=BODY_FONT,
                         italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if rest:
                add_text(slide, Inches(1.2), Inches(4.95), Inches(10.93), Inches(1.2),
                         "   |   ".join(rest), 16,
                         RGBColor(0xC9, 0xB8, 0x8C), font=BODY_FONT,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            if s["visual"]:
                add_text(slide, Inches(1.2), Inches(6.25), Inches(10.93), Inches(0.6),
                         "Visuel suggere : " + s["visual"], 10,
                         RGBColor(0x9A, 0x88, 0x9A), font=BODY_FONT,
                         italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            # Slide de contenu standard
            add_text(slide, Inches(0.5), Inches(0.78), Inches(12.33), Inches(1.0),
                     s["title"], 34, title_color, font=TITLE_FONT, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_rect(slide, Inches(0.55), Inches(1.72), Inches(2.4), Emu(34000), GOLD)

            body_color = ANTHRACITE
            col_x = Inches(0.55)
            col_y = Inches(2.05)
            col_w = Inches(6.6)
            col_h = Inches(4.4)

            bullets = s["bullets"]
            if bullets:
                nb = len(bullets)
                bsize = 20 if nb <= 4 else (18 if nb <= 6 else 16)
                add_bullets(slide, col_x, col_y, col_w, col_h,
                            bullets, bsize, body_color)

            # Colonne droite : espace visuel reserve + suggestion d'icone
            add_visual_placeholder(slide, Inches(7.45), Inches(2.05),
                                   Inches(5.35), Inches(4.2), s["visual"], on_dark)

            if s["todo"]:
                add_todo_badge(slide, Inches(0.55), Inches(6.45))

        add_footer(slide, on_dark)
        add_pagenum(slide, str(s["n"]) + " / " + str(total), on_dark)
        add_notes(slide, s["notes"])

    # --- Intercalaire backup ---
    sep = prs.slides.add_slide(blank)
    set_bg(sep, AUBERGINE)
    add_rect(sep, 0, 0, SLIDE_W, Inches(0.22), AUBERGINE)
    add_rect(sep, 0, Inches(0.22), SLIDE_W, Emu(28000), GOLD)
    add_text(sep, Inches(1.0), Inches(2.6), Inches(11.33), Inches(1.4),
             "Slides annexes (backup)", 44, GOLD, font=TITLE_FONT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(sep, Inches(1.2), Inches(4.1), Inches(10.93), Inches(1.0),
             "A ne pas presenter - uniquement pour repondre aux questions du jury",
             20, LIGHT_ON_DARK, font=BODY_FONT, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(sep, True)
    add_notes(sep, "Intercalaire : les slides suivantes sont des annexes de secours, "
                   "preparees mais non presentees, pour repondre precisement aux "
                   "questions du jury sans surcharger la presentation principale.")

    # --- Slides annexes ---
    for b in backups:
        slide = prs.slides.add_slide(blank)
        set_bg(slide, IVORY)
        add_rect(slide, 0, 0, SLIDE_W, Inches(0.22), GOLD)
        add_rect(slide, 0, Inches(0.22), SLIDE_W, Emu(28000), PLUM)
        add_text(slide, Inches(0.5), Inches(0.32), Inches(8), Inches(0.4),
                 "ANNEXE (BACKUP)", 12, GOLD, font=BODY_FONT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(0.5), Inches(0.78), Inches(12.33), Inches(1.0),
                 b["title"], 28, PLUM, font=TITLE_FONT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(0.55), Inches(1.72), Inches(2.4), Emu(34000), GOLD)
        if b["bullets"]:
            add_bullets(slide, Inches(0.55), Inches(2.05), Inches(6.6), Inches(4.4),
                        b["bullets"], 18, ANTHRACITE)
        add_visual_placeholder(slide, Inches(7.45), Inches(2.05),
                               Inches(5.35), Inches(4.2), b["visual"], False)
        add_footer(slide, False)
        add_pagenum(slide, b["code"], False)
        add_notes(slide, b["notes"])

    prs.save(OUTPUT)
    return OUTPUT, len(prs.slides._sldIdLst), total, len(backups)


if __name__ == "__main__":
    path, count, main_count, backup_count = build()
    print("OK -> " + path)
    print("Slides principales parsees : " + str(main_count))
    print("Slides annexes (backup)    : " + str(backup_count))
    print("Slides generees au total   : " + str(count))
